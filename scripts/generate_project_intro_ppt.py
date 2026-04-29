from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("artifacts") / "AI_Document_Review_System_Intro.pptx"


BG = RGBColor(244, 248, 251)
SURFACE = RGBColor(255, 255, 255)
SURFACE_ALT = RGBColor(248, 251, 253)
BORDER = RGBColor(222, 234, 241)
TEXT = RGBColor(23, 43, 77)
TEXT_SOFT = RGBColor(91, 111, 130)
MUTED = RGBColor(132, 150, 169)
BLUE = RGBColor(47, 128, 237)
TEAL = RGBColor(44, 143, 153)
GREEN = RGBColor(39, 174, 96)
INDIGO = RGBColor(92, 99, 200)
ORANGE = RGBColor(242, 153, 74)
SOFT_BLUE = RGBColor(235, 244, 254)
SOFT_TEAL = RGBColor(231, 245, 246)
SOFT_GREEN = RGBColor(237, 248, 239)
SOFT_ORANGE = RGBColor(255, 245, 234)
SOFT_INDIGO = RGBColor(239, 240, 251)


def set_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text="", font_size=18, bold=False, color=TEXT,
                font_name="Yu Gothic", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_card(slide, left, top, width, height, fill=SURFACE, line=BORDER, radius=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.0)
    return shape


def add_pill(slide, left, top, width, height, text, fill=BLUE, color=SURFACE, font_size=15):
    pill = add_card(slide, left, top, width, height, fill=fill, line=fill)
    tf = pill.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Yu Gothic"
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = color
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return pill


def add_metric_card(slide, left, top, width, height, accent, soft_fill, value, label):
    card = add_card(slide, left, top, width, height)
    # icon container
    icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + Inches(0.2), top + Inches(0.18), Inches(0.7), Inches(0.55))
    icon.fill.solid()
    icon.fill.fore_color.rgb = soft_fill
    icon.line.color.rgb = soft_fill
    tf = icon.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "◻"
    run.font.name = "Yu Gothic"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = accent
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_textbox(slide, left + Inches(0.2), top + Inches(0.95), width - Inches(0.4), Inches(0.8), value, 34, True, accent, align=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.22), top + Inches(1.55), width - Inches(0.44), Inches(0.7), label, 15, True, TEXT, align=PP_ALIGN.CENTER)
    accent_line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left + Inches(0.35), top + height - Inches(0.08), width - Inches(0.7), Inches(0.04))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = accent
    accent_line.line.color.rgb = accent
    return card


def add_section_title(slide, eyebrow, title, subtitle=None):
    add_textbox(slide, Inches(0.7), Inches(0.35), Inches(4.5), Inches(0.35), eyebrow, 12, True, TEAL)
    add_textbox(slide, Inches(0.7), Inches(0.72), Inches(7.6), Inches(0.9), title, 26, True, TEXT)
    if subtitle:
        add_textbox(slide, Inches(0.7), Inches(1.55), Inches(8.5), Inches(0.45), subtitle, 12.5, False, TEXT_SOFT)


def add_bullet_list(slide, left, top, width, items, color=TEXT_SOFT, bullet_color=TEAL, font_size=18):
    y = top
    for item in items:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, y + Inches(0.08), Inches(0.09), Inches(0.09))
        dot.fill.solid()
        dot.fill.fore_color.rgb = bullet_color
        dot.line.color.rgb = bullet_color
        add_textbox(slide, left + Inches(0.18), y, width - Inches(0.18), Inches(0.42), item, font_size, False, color)
        y += Inches(0.48)


def add_small_kpi(slide, left, top, width, title, value, accent):
    card = add_card(slide, left, top, width, Inches(1.25), SURFACE, BORDER)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.14), width - Inches(0.36), Inches(0.3), title, 11, True, TEXT_SOFT)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.44), width - Inches(0.36), Inches(0.55), value, 24, True, accent)
    return card


def add_abstract_ai_illustration(slide, left, top):
    base = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + Inches(0.65), top + Inches(1.55), Inches(1.9), Inches(0.38))
    base.fill.solid()
    base.fill.fore_color.rgb = SOFT_BLUE
    base.line.color.rgb = SOFT_BLUE

    platform = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + Inches(0.35), top + Inches(1.15), Inches(2.5), Inches(0.7))
    platform.fill.solid()
    platform.fill.fore_color.rgb = SURFACE
    platform.line.color.rgb = BORDER

    brain = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(1.1), top + Inches(0.2), Inches(1.1), Inches(1.1))
    brain.fill.solid()
    brain.fill.fore_color.rgb = SOFT_BLUE
    brain.line.color.rgb = BLUE
    brain.line.width = Pt(1.8)

    glow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.95), top + Inches(0.05), Inches(1.4), Inches(1.4))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(232, 243, 255)
    glow.line.color.rgb = RGBColor(232, 243, 255)
    glow.shadow.inherit = False

    # reorder glow behind brain by creating brain after glow would be easier; acceptable
    panel1 = add_card(slide, left + Inches(0.1), top + Inches(0.5), Inches(0.95), Inches(0.75), fill=SURFACE_ALT)
    panel2 = add_card(slide, left + Inches(2.35), top + Inches(0.7), Inches(0.95), Inches(0.75), fill=SURFACE_ALT)
    panel3 = add_card(slide, left + Inches(2.05), top + Inches(1.45), Inches(1.1), Inches(0.78), fill=SURFACE_ALT)
    for panel, accent in [(panel1, BLUE), (panel2, TEAL), (panel3, GREEN)]:
        for i in range(3):
            line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, panel.left + Inches(0.12), panel.top + Inches(0.14 + 0.12 * i), Inches(0.55), Inches(0.03))
            line.fill.solid()
            line.fill.fore_color.rgb = accent
            line.line.color.rgb = accent


def add_process_step(slide, left, top, width, number, title, subtitle, accent, soft_fill):
    card = add_card(slide, left, top, width, Inches(2.1))
    badge = add_pill(slide, left + Inches(0.12), top + Inches(0.1), Inches(0.55), Inches(0.35), number, fill=accent, font_size=14)
    del badge
    add_textbox(slide, left + Inches(0.76), top + Inches(0.1), width - Inches(0.88), Inches(0.38), title, 15, True, accent)
    icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + Inches(0.18), top + Inches(0.62), Inches(0.72), Inches(0.58))
    icon.fill.solid()
    icon.fill.fore_color.rgb = soft_fill
    icon.line.color.rgb = soft_fill
    add_textbox(slide, left + Inches(0.18), top + Inches(0.7), Inches(0.72), Inches(0.25), "▣", 18, True, accent, align=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.18), top + Inches(1.32), width - Inches(0.36), Inches(0.58), subtitle, 12.5, False, TEXT_SOFT)
    return card


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1 cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    hero = add_card(slide, Inches(0.35), Inches(0.35), Inches(12.65), Inches(6.75), fill=SURFACE)
    hero.shadow.inherit = False

    add_textbox(slide, Inches(0.9), Inches(0.7), Inches(5.5), Inches(0.35), "AI DOCUMENT REVIEW PLATFORM", 12, True, TEAL)
    add_textbox(slide, Inches(0.9), Inches(1.1), Inches(6.3), Inches(1.0), "Hệ thống AI Review Tài Liệu", 27, True, TEXT)
    add_textbox(slide, Inches(0.9), Inches(2.0), Inches(4.8), Inches(0.5), "Nền tảng review tài liệu đa loại bằng AI", 17, True, SURFACE, align=PP_ALIGN.CENTER)
    pill_bg = add_pill(slide, Inches(0.9), Inches(2.0), Inches(4.8), Inches(0.52), "Nền tảng review tài liệu đa loại bằng AI", fill=BLUE, font_size=16)
    pill_bg.shadow.inherit = False
    add_abstract_ai_illustration(slide, Inches(8.2), Inches(0.55))
    add_metric_card(slide, Inches(0.95), Inches(3.55), Inches(2.45), Inches(2.35), BLUE, SOFT_BLUE, "4", "Loại tài liệu review")
    add_metric_card(slide, Inches(3.8), Inches(3.55), Inches(2.45), Inches(2.35), ORANGE, SOFT_ORANGE, "2", "Ngôn ngữ hỗ trợ")
    add_metric_card(slide, Inches(6.65), Inches(3.55), Inches(2.45), Inches(2.35), GREEN, SOFT_GREEN, "v1", "Rubric version active")
    add_metric_card(slide, Inches(9.5), Inches(3.55), Inches(2.45), Inches(2.35), INDIGO, SOFT_INDIGO, "AI", "Gemini grading")

    # Slide 2 problem
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_section_title(slide, "01 / Problem", "Bài toán hiện tại", "Những điểm nghẽn khi review tài liệu thủ công trong dự án phần mềm")
    add_card(slide, Inches(0.7), Inches(1.8), Inches(5.9), Inches(4.95))
    add_card(slide, Inches(6.85), Inches(1.8), Inches(5.8), Inches(4.95))
    add_textbox(slide, Inches(1.0), Inches(2.1), Inches(2.8), Inches(0.4), "Thách thức", 18, True, TEXT)
    add_bullet_list(
        slide,
        Inches(1.0),
        Inches(2.55),
        Inches(5.1),
        [
            "Review thủ công tốn thời gian và khó giữ chất lượng ổn định",
            "Tiêu chí đánh giá dễ lệch giữa reviewer và giữa các dự án",
            "Khó mở rộng khi phải hỗ trợ nhiều loại tài liệu khác nhau",
            "Feedback thường dài, khó chuẩn hóa và khó tái sử dụng",
        ],
    )
    add_textbox(slide, Inches(7.15), Inches(2.1), Inches(3.0), Inches(0.4), "Kỳ vọng từ hệ thống", 18, True, TEXT)
    add_small_kpi(slide, Inches(7.15), Inches(2.65), Inches(1.8), "Chuẩn hóa", "Same rubric", TEAL)
    add_small_kpi(slide, Inches(9.15), Inches(2.65), Inches(1.6), "Tốc độ", "Faster", BLUE)
    add_small_kpi(slide, Inches(10.95), Inches(2.65), Inches(1.45), "Mở rộng", "Multi-type", GREEN)
    add_bullet_list(
        slide,
        Inches(7.15),
        Inches(4.0),
        Inches(5.0),
        [
            "Chọn đúng loại tài liệu và dùng rubric tương ứng",
            "Tự động trả tổng điểm, điểm theo tiêu chí và gợi ý cải thiện",
            "Dễ cập nhật rubric theo version mà không phá hệ thống",
        ],
        bullet_color=BLUE,
    )

    # Slide 3 solution/process
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_section_title(slide, "02 / Solution", "Luồng giải pháp tổng thể", "Từ upload tài liệu đến kết quả review và feedback có cấu trúc")
    step_top = Inches(2.0)
    step_left = Inches(0.7)
    step_w = Inches(2.35)
    gap = Inches(0.23)
    steps = [
        ("01", "Upload", "PDF / PPTX", TEAL, SOFT_TEAL),
        ("02", "Classify", "document_type", BLUE, SOFT_BLUE),
        ("03", "Resolve rubric", "language + version", INDIGO, SOFT_INDIGO),
        ("04", "AI grading", "Gemini scoring", TEAL, SOFT_TEAL),
        ("05", "Result", "score + feedback", GREEN, SOFT_GREEN),
    ]
    for idx, step in enumerate(steps):
        add_process_step(slide, step_left + idx * (step_w + gap), step_top, step_w, *step)
        if idx < len(steps) - 1:
            cx = step_left + idx * (step_w + gap) + step_w + Inches(0.04)
            cy = step_top + Inches(1.05)
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, cx, cy, cx + Inches(0.13), cy)
            conn.line.color.rgb = MUTED
            conn.line.width = Pt(2.0)
    add_card(slide, Inches(0.7), Inches(5.0), Inches(12.0), Inches(1.2), fill=SURFACE_ALT)
    add_bullet_list(
        slide,
        Inches(1.0),
        Inches(5.3),
        Inches(11.2),
        [
            "Frontend cho phép chọn loại tài liệu trước khi upload",
            "Backend chọn rubric theo document_type + language + version",
            "Kết quả gồm score, criteria_scores, criteria_suggestions và draft_feedback",
        ],
        bullet_color=GREEN,
        font_size=16,
    )

    # Slide 4 document types
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_section_title(slide, "03 / Document Types", "4 loại tài liệu đang hỗ trợ")
    types = [
        ("Project Review", "Nhìn nhận dự án,\nretrospective,\nbài học và chính sách", TEAL, SOFT_TEAL),
        ("Bug Analysis", "Phân tích bug,\nnguyên nhân,\nảnh hưởng và tái phát", BLUE, SOFT_BLUE),
        ("QA Review", "Chất lượng kiểm thử,\nđộ bao phủ,\ndefect và quy trình", GREEN, SOFT_GREEN),
        ("Explanation Review", "Tính rõ ràng,\nlogic truyền đạt\nvà khả năng giải thích", INDIGO, SOFT_INDIGO),
    ]
    for idx, (title, body, accent, soft) in enumerate(types):
        left = Inches(0.75 + idx * 3.1)
        card = add_card(slide, left, Inches(2.0), Inches(2.75), Inches(3.7))
        icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + Inches(0.2), Inches(2.22), Inches(0.78), Inches(0.62))
        icon.fill.solid()
        icon.fill.fore_color.rgb = soft
        icon.line.color.rgb = soft
        add_textbox(slide, left + Inches(0.2), Inches(2.34), Inches(0.78), Inches(0.24), "▣", 20, True, accent, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.2), Inches(3.0), Inches(2.25), Inches(0.5), title, 17, True, accent)
        add_textbox(slide, left + Inches(0.2), Inches(3.55), Inches(2.3), Inches(1.4), body, 13, False, TEXT_SOFT)

    # Slide 5 architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_section_title(slide, "04 / Architecture", "Kiến trúc hệ thống")
    blocks = [
        (Inches(0.95), Inches(2.2), Inches(2.45), Inches(2.25), "Frontend", "React + TypeScript + Vite\nUpload, dashboard, results", BLUE, SOFT_BLUE),
        (Inches(3.8), Inches(2.2), Inches(2.45), Inches(2.25), "Backend", "FastAPI\nupload / grading / jobs", TEAL, SOFT_TEAL),
        (Inches(6.65), Inches(2.2), Inches(2.45), Inches(2.25), "Rubric", "document_type + version + language", INDIGO, SOFT_INDIGO),
        (Inches(9.5), Inches(2.2), Inches(2.45), Inches(2.25), "AI", "Gemini\nJSON structured result", GREEN, SOFT_GREEN),
    ]
    for left, top, width, height, title, body, accent, soft in blocks:
        add_card(slide, left, top, width, height)
        chip = add_pill(slide, left + Inches(0.18), top + Inches(0.16), Inches(1.05), Inches(0.34), title, fill=accent, font_size=12)
        del chip
        add_textbox(slide, left + Inches(0.2), top + Inches(0.8), width - Inches(0.4), Inches(0.9), body, 15, False, TEXT)
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left + Inches(0.2), top + height - Inches(0.08), width - Inches(0.4), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = accent
        line.line.color.rgb = accent
    for x in [Inches(3.45), Inches(6.3), Inches(9.15)]:
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, Inches(3.3), x + Inches(0.18), Inches(3.3))
        conn.line.color.rgb = MUTED
        conn.line.width = Pt(2.2)

    # Slide 6 result structure
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_section_title(slide, "05 / Output", "Cấu trúc kết quả review")
    add_card(slide, Inches(0.8), Inches(1.95), Inches(4.0), Inches(4.7))
    add_textbox(slide, Inches(1.05), Inches(2.2), Inches(2.0), Inches(0.35), "Project Review schema", 18, True, TEXT)
    add_bullet_list(
        slide,
        Inches(1.05),
        Inches(2.7),
        Inches(3.4),
        [
            "review_tong_the /25",
            "diem_tot /25",
            "diem_xau /30",
            "chinh_sach /20",
        ],
        bullet_color=TEAL,
        font_size=16,
    )
    add_card(slide, Inches(5.15), Inches(1.95), Inches(7.0), Inches(4.7))
    add_textbox(slide, Inches(5.4), Inches(2.2), Inches(2.8), Inches(0.35), "Response payload", 18, True, TEXT)
    code_lines = [
        '{',
        '  "score": 88,',
        '  "criteria_scores": {...},',
        '  "criteria_suggestions": {...},',
        '  "draft_feedback": "..."',
        '}',
    ]
    add_textbox(slide, Inches(5.5), Inches(2.75), Inches(6.1), Inches(2.2), "\n".join(code_lines), 17, False, TEXT)
    add_bullet_list(
        slide,
        Inches(5.45),
        Inches(5.25),
        Inches(6.2),
        [
            "criteria_suggestions giúp người dùng hiểu lý do bị trừ điểm",
            "Feedback tổng hợp hiển thị theo ngôn ngữ của bài review",
        ],
        bullet_color=BLUE,
        font_size=15,
    )

    # Slide 7 strengths
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_section_title(slide, "06 / Strengths", "Điểm mạnh nổi bật")
    strengths = [
        ("Versioned rubric", "Cập nhật rubric theo version, không phá flow hiện tại", TEAL, SOFT_TEAL),
        ("Multi-language", "Rubric và UI hỗ trợ tiếng Việt / tiếng Nhật", BLUE, SOFT_BLUE),
        ("Batch grading", "Có background job cho chấm hàng loạt", GREEN, SOFT_GREEN),
        ("Signature-safe regrade", "Kết quả gắn với content_hash, rubric_version, gemini_model", INDIGO, SOFT_INDIGO),
        ("Detailed guidance", "Hiển thị giải thích và hướng tăng điểm theo từng tiêu chí", ORANGE, SOFT_ORANGE),
        ("Extensible types", "Hỗ trợ nhiều document_type trong cùng một nền tảng", TEAL, SOFT_TEAL),
    ]
    for idx, (title, body, accent, soft) in enumerate(strengths):
        col = idx % 3
        row = idx // 3
        left = Inches(0.75 + col * 4.05)
        top = Inches(2.0 + row * 2.05)
        add_card(slide, left, top, Inches(3.65), Inches(1.7))
        icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + Inches(0.18), top + Inches(0.18), Inches(0.6), Inches(0.5))
        icon.fill.solid()
        icon.fill.fore_color.rgb = soft
        icon.line.color.rgb = soft
        add_textbox(slide, left + Inches(0.95), top + Inches(0.18), Inches(2.2), Inches(0.3), title, 15, True, accent)
        add_textbox(slide, left + Inches(0.95), top + Inches(0.6), Inches(2.4), Inches(0.6), body, 12.5, False, TEXT_SOFT)

    # Slide 8 roadmap
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_section_title(slide, "07 / Roadmap", "Hướng phát triển tiếp theo", "Những hạng mục có thể tăng giá trị vận hành và quản trị hệ thống")
    roadmap = [
        ("Regrade all", "Chấm lại toàn bộ khi đổi rubric version", TEAL),
        ("Audit trail", "Theo dõi rõ rubric_version cho từng submission", BLUE),
        ("Report export", "Xuất báo cáo tổng hợp theo dự án / loại tài liệu", GREEN),
        ("Rubric workflow", "Quản lý version và promote rubric tốt hơn", INDIGO),
        ("Admin tools", "Cấu hình active version và kiểm tra metadata", ORANGE),
    ]
    for idx, (title, body, accent) in enumerate(roadmap):
        left = Inches(1.0 + idx * 2.35)
        add_pill(slide, left, Inches(2.2), Inches(1.55), Inches(0.42), f"{idx + 1:02d}", fill=accent, font_size=14)
        add_card(slide, left - Inches(0.1), Inches(2.75), Inches(1.75), Inches(2.2))
        add_textbox(slide, left + Inches(0.05), Inches(3.05), Inches(1.35), Inches(0.4), title, 14, True, accent, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.04), Inches(3.55), Inches(1.38), Inches(0.85), body, 11.5, False, TEXT_SOFT, align=PP_ALIGN.CENTER)
        if idx < len(roadmap) - 1:
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left + Inches(1.7), Inches(3.85), left + Inches(2.05), Inches(3.85))
            conn.line.color.rgb = MUTED
            conn.line.width = Pt(2.0)

    return prs


def main():
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(str(OUT_PATH))
    print(f"Saved {OUT_PATH}")


if __name__ == "__main__":
    main()
