import re
import pdfplumber
from pathlib import Path


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF using pdfplumber for better multilingual support."""
    text_parts = []
    try:
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"[Page {page_num}]\n{page_text}")
    except Exception as e:
        print(f"[PDF Extract] Error with pdfplumber: {e}")
        try:
            # Fallback to PyPDF2
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            for page_num, page in enumerate(reader.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"[Page {page_num}]\n{page_text}")
        except Exception as e2:
            print(f"[PDF Extract] Error with PyPDF2: {e2}")
    return "\n\n".join(text_parts)


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PowerPoint (.pptx) file."""
    from pptx import Presentation
    
    text_parts = []
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
    except Exception as e:
        print(f"[PPTX Extract] Error extracting text: {e}")
        return ""
    
    return "\n\n".join(text_parts)


def extract_text_from_file(file_path: str) -> str:
    """Extract text from either PDF or PowerPoint file."""
    file_path = Path(file_path)
    extension = file_path.suffix.lower()
    
    if extension == '.pdf':
        return extract_text_from_pdf(str(file_path))
    elif extension == '.pptx':
        return extract_text_from_pptx(str(file_path))
    else:
        raise ValueError(f"Unsupported file format: {extension}. Only PDF and PowerPoint (.pptx) are supported.")


def detect_language_from_text(text: str) -> str:
    """
    Detect language from text content using weighted scoring.
    Returns: 'vi' for Vietnamese, 'ja' for Japanese
    
    Strategy: Use weighted score to handle mixed-language documents
    """
    text_length = len(text)
    if text_length == 0:
        return "ja"  # Default to Japanese
    
    # Count characteristic characters for each language
    # Japanese: Hiragana, Katakana, Kanji
    ja_chars = len(re.findall(r'[\u3040-\u309F\u30A0-\u30FF\u4E00-\u9FAF]', text))
    
    # Vietnamese: Characters with diacritics
    vi_chars = len(re.findall(r'[àáạảãâầấậẩẫăằắặẳẵèéẹẻẽêềếệểễìíịỉĩòóọỏõôồốộổỗơờớợởỡùúụủũưừứựửữỳýỵỷỹđÀÁẠẢÃÂẦẤẬẨẪĂẰẮẶẲẴÈÉẸẺẼÊỀẾỆỂỄÌÍỊỈĨÒÓỌỎÕÔỒỐỘỔỖƠỜỚỢỞỠÙÚỤỦŨƯỪỨỰỬỮỲÝỴỶỸĐ]', text))
    
    # English: Basic Latin characters
    en_chars = len(re.findall(r'[a-zA-Z]', text))
    
    # Japanese-specific patterns (common words/particles)
    ja_patterns = len(re.findall(r'\b(です|ます|した|ている|こと|もの|これ|それ|あれ|この|その|あの|は|が|を|に|へ|で|と|から|まで|より|も|や|など|か|ね|よ)\b', text))
    
    # Vietnamese-specific patterns
    vi_patterns = len(re.findall(r'\b(của|là|các|và|cho|đã|được|không|này|nhưng|với|một|người|tôi|chúng|có|về|những|tại|để)\b', text, re.IGNORECASE))
    
    # Calculate ratios (0-1)
    ja_ratio = ja_chars / text_length if text_length > 0 else 0
    vi_ratio = vi_chars / text_length if text_length > 0 else 0
    
    # Calculate weighted scores
    # Weight factors: characters (40%), patterns (30%), ratio (30%)
    ja_score = (ja_chars * 0.4) + (ja_patterns * 10 * 0.3) + (ja_ratio * 1000 * 0.3)
    vi_score = (vi_chars * 0.4) + (vi_patterns * 10 * 0.3) + (vi_ratio * 1000 * 0.3)
    
    # Log detection details for debugging
    print(f"[Language Detection] JA chars: {ja_chars} ({ja_ratio:.4f}), VI chars: {vi_chars} ({vi_ratio:.4f})")
    print(f"[Language Detection] JA patterns: {ja_patterns}, VI patterns: {vi_patterns}")
    print(f"[Language Detection] JA score: {ja_score:.2f}, VI score: {vi_score:.2f}")
    
    # Determine language based on weighted score
    # Require significant difference to avoid false positives from mixed content
    if ja_score > vi_score and ja_score > 50:  # Japanese dominant
        detected = "ja"
    elif vi_score > ja_score and vi_score > 50:  # Vietnamese dominant
        detected = "vi"
    elif ja_chars > 0 and vi_chars > 0:  # Mixed content - use ratio
        detected = "ja" if ja_ratio > vi_ratio else "vi"
    elif ja_chars > 10:  # Fallback: significant Japanese content
        detected = "ja"
    elif vi_chars > 10:  # Fallback: significant Vietnamese content
        detected = "vi"
    else:
        detected = "ja"  # Default to Japanese
    
    print(f"[Language Detection] Detected: {detected}")
    return detected
